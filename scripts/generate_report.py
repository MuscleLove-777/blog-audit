#!/usr/bin/env python3
"""
監査結果パワポレポート生成スクリプト
results/latest.json を読み込み、.pptx形式で監査レポートを生成
"""

import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 色定義
BG_COLOR = RGBColor(0x0B, 0x14, 0x2E)  # 濃紺
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0xA0, 0x43)
YELLOW = RGBColor(0xD2, 0x99, 0x22)
RED = RGBColor(0xF8, 0x51, 0x49)
BLUE = RGBColor(0x58, 0xA6, 0xFF)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_BG = RGBColor(0x16, 0x1B, 0x22)


def set_slide_bg(slide, color=BG_COLOR):
    """スライド背景色を設定"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    """図形を追加"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_status_circle(slide, left, top, size, status):
    """ステータス信号機を追加"""
    color = {"green": GREEN, "yellow": YELLOW, "red": RED}.get(status, GRAY)
    shape = add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill_color=color)
    return shape


def slide_title(slide, title_text, subtitle_text=None):
    """スライドタイトルを図解的に追加"""
    # タイトル下線
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(3), Emu(40000), fill_color=BLUE)
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8), title_text, font_size=28, color=WHITE, bold=True)
    if subtitle_text:
        add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5), subtitle_text, font_size=14, color=GRAY)


def create_cover_slide(prs, timestamp):
    """スライド1: 表紙"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
    set_slide_bg(slide)

    # 装飾図形（グラデーション風の四角形）
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15), fill_color=BLUE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), Inches(10), Inches(0.15), fill_color=BLUE)

    # 左サイドの装飾線
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Emu(40000), Inches(3.5), fill_color=BLUE)

    # タイトル
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(8), Inches(1.0),
                 "ブログ監査レポート", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    # 日付
    date_str = timestamp.split(" ")[0] if " " in timestamp else timestamp
    add_text_box(slide, Inches(1.2), Inches(3.3), Inches(8), Inches(0.6),
                 date_str, font_size=24, color=BLUE, bold=False, alignment=PP_ALIGN.LEFT)

    # サブタイトル
    add_text_box(slide, Inches(1.2), Inches(4.2), Inches(8), Inches(0.6),
                 "全9サイト自動監査結果", font_size=18, color=GRAY)

    # アイコン的な図形群
    for i in range(3):
        add_shape(slide, MSO_SHAPE.OVAL, Inches(7.5 + i * 0.6), Inches(5.5), Inches(0.4), Inches(0.4),
                  fill_color=[GREEN, YELLOW, RED][i])


def create_summary_slide(prs, sites):
    """スライド2: 全体サマリー（信号機表示）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_title(slide, "全体サマリー", "各サイトの健康度")

    # 信号機凡例
    for i, (label, color) in enumerate([("正常", GREEN), ("警告", YELLOW), ("エラー", RED)]):
        x = Inches(6.5 + i * 1.2)
        add_shape(slide, MSO_SHAPE.OVAL, x, Inches(0.5), Inches(0.3), Inches(0.3), fill_color=color)
        add_text_box(slide, x + Inches(0.35), Inches(0.45), Inches(0.8), Inches(0.4), label, font_size=10, color=GRAY)

    # サイトカード（3列グリッド）
    cols = 3
    card_w = Inches(2.8)
    card_h = Inches(1.4)
    start_x = Inches(0.5)
    start_y = Inches(1.8)
    gap_x = Inches(3.0)
    gap_y = Inches(1.6)

    for idx, site in enumerate(sites):
        col = idx % cols
        row = idx // cols
        x = start_x + col * gap_x
        y = start_y + row * gap_y

        # カード背景
        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, fill_color=LIGHT_BG, line_color=GRAY)

        # ステータス信号
        status_color = {"green": GREEN, "yellow": YELLOW, "red": RED}[site["status"]]
        add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.35), Inches(0.35), fill_color=status_color)

        # サイト名
        add_text_box(slide, x + Inches(0.6), y + Inches(0.1), Inches(2.0), Inches(0.4),
                     site["name"], font_size=11, color=WHITE, bold=True)

        # ジャンル
        add_text_box(slide, x + Inches(0.6), y + Inches(0.45), Inches(1.5), Inches(0.3),
                     f"[{site['genre']}]", font_size=9, color=GRAY)

        # スコア
        score = site["score"]
        add_text_box(slide, x + Inches(0.15), y + Inches(0.85), Inches(2.5), Inches(0.4),
                     f"スコア: {score['passed']}/{score['total']} ({score['percentage']}%)",
                     font_size=10, color=WHITE)


def create_detail_slide(prs, sites):
    """スライド3: 各サイト詳細テーブル"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_title(slide, "各サイト詳細チェック結果", "全チェック項目一覧")

    # テーブル作成
    rows = len(sites) + 1
    cols = 7
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.6), Inches(9.4), Inches(5.0))
    table = table_shape.table

    # ヘッダー
    headers = ["サイト", "HTTP", "CSS", "記事数", "リンク", "SEO", "スコア"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x2A, 0x40)

    # データ行
    for row_idx, site in enumerate(sites):
        checks = site["checks"]
        row_data = [
            f"{site['name']}\n({site['genre']})",
            "OK" if checks["http"]["ok"] else "NG",
            "OK" if checks["css"]["ok"] else "NG",
            str(checks["articles"].get("count", 0)),
            "OK" if checks["affiliate_links"]["ok"] else "NG",
            "OK" if checks["sitemap"]["ok"] and checks["robots_txt"]["ok"] else "NG",
            f"{site['score']['passed']}/{site['score']['total']}",
        ]
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = WHITE
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if row_idx % 2 == 0 else BG_COLOR

    # 図解: 右下にステータス分布の棒グラフ風
    green_count = sum(1 for s in sites if s["status"] == "green")
    yellow_count = sum(1 for s in sites if s["status"] == "yellow")
    red_count = sum(1 for s in sites if s["status"] == "red")
    total = len(sites)

    # ミニ棒グラフ
    bar_base_x = Inches(7.5)
    bar_base_y = Inches(1.6)
    if total > 0:
        for i, (count, color, label) in enumerate([
            (green_count, GREEN, "正常"),
            (yellow_count, YELLOW, "警告"),
            (red_count, RED, "エラー"),
        ]):
            bar_height = Inches(count / total * 1.2) if count > 0 else Emu(10000)
            y_pos = bar_base_y + Inches(1.2) - bar_height + Inches(i * 0)
            # テキストラベルのみ
            pass


def create_seo_slide(prs, sites):
    """スライド4: SEOチェック結果"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_title(slide, "SEOチェック結果", "サイトマップ / robots.txt / メタタグ")

    # SEO項目ごとに図解
    items = ["sitemap", "robots_txt", "responsive", "dark_theme"]
    item_labels = {
        "sitemap": "サイトマップ",
        "robots_txt": "robots.txt",
        "responsive": "レスポンシブ対応",
        "dark_theme": "ダークテーマ",
    }

    start_y = Inches(1.8)
    for col_idx, item in enumerate(items):
        x = Inches(0.3 + col_idx * 2.4)
        # ヘッダーカード
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, Inches(2.2), Inches(0.5), fill_color=LIGHT_BG, line_color=BLUE)
        add_text_box(slide, x + Inches(0.1), start_y + Inches(0.05), Inches(2.0), Inches(0.4),
                     item_labels[item], font_size=11, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

        for row_idx, site in enumerate(sites):
            y = start_y + Inches(0.7) + row_idx * Inches(0.5)
            ok = site["checks"].get(item, {}).get("ok", False)
            status_color = GREEN if ok else RED
            icon_text = "OK" if ok else "NG"

            # 信号アイコン
            add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.05), y + Inches(0.05), Inches(0.25), Inches(0.25), fill_color=status_color)
            add_text_box(slide, x + Inches(0.4), y, Inches(1.8), Inches(0.35),
                         f"{site['name'][:12]}", font_size=8, color=WHITE)


def create_issues_slide(prs, sites):
    """スライド5: 問題点と推奨アクション"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    slide_title(slide, "問題点と推奨アクション", "改善が必要な項目")

    # 問題点を収集
    issues = []
    for site in sites:
        checks = site["checks"]
        name = site["name"]
        if not checks["http"]["ok"]:
            issues.append({"site": name, "issue": "サイトがダウン中", "severity": "red", "action": "デプロイ状況を確認"})
        if not checks["css"]["ok"]:
            issues.append({"site": name, "issue": "CSSが読み込まれていない", "severity": "red", "action": "stylesheet参照を修正"})
        if not checks["articles"]["ok"]:
            issues.append({"site": name, "issue": "記事が0件", "severity": "red", "action": "記事自動生成の確認"})
        if not checks["affiliate_links"]["ok"]:
            issues.append({"site": name, "issue": "アフィリリンクなし", "severity": "yellow", "action": "アフィリリンクを追加"})
        if not checks["musclelove_links"]["ok"]:
            issues.append({"site": name, "issue": "SNSリンクなし", "severity": "yellow", "action": "Patreon/X/Linktreeリンクを追加"})
        if not checks["sitemap"]["ok"]:
            issues.append({"site": name, "issue": "サイトマップなし", "severity": "yellow", "action": "sitemap.xmlを追加"})
        if not checks["robots_txt"]["ok"]:
            issues.append({"site": name, "issue": "robots.txtなし", "severity": "yellow", "action": "robots.txtを追加"})
        if not checks["responsive"]["ok"]:
            issues.append({"site": name, "issue": "viewportメタタグなし", "severity": "yellow", "action": "metaタグを追加"})

    if not issues:
        # 問題なしの場合
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(3), Inches(6), Inches(1.5), fill_color=LIGHT_BG, line_color=GREEN)
        add_text_box(slide, Inches(2.5), Inches(3.3), Inches(5), Inches(0.8),
                     "問題は検出されませんでした", font_size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
        return

    # 問題リスト表示（左: 問題 / 右: アクション）
    # 左パネル
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.7), Inches(4.5), Inches(0.5), fill_color=RED, line_color=RED)
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(4.0), Inches(0.4),
                 "検出された問題", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 右パネル
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.7), Inches(4.5), Inches(0.5), fill_color=GREEN, line_color=GREEN)
    add_text_box(slide, Inches(5.4), Inches(1.75), Inches(4.0), Inches(0.4),
                 "推奨アクション", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 矢印
    add_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(4.6), Inches(3.5), Inches(0.7), Inches(0.4), fill_color=BLUE)

    max_display = min(len(issues), 10)
    for i in range(max_display):
        issue = issues[i]
        y = Inches(2.4) + i * Inches(0.45)
        severity_color = RED if issue["severity"] == "red" else YELLOW

        # 問題（左）
        add_shape(slide, MSO_SHAPE.OVAL, Inches(0.4), y + Inches(0.05), Inches(0.2), Inches(0.2), fill_color=severity_color)
        add_text_box(slide, Inches(0.7), y, Inches(4.0), Inches(0.4),
                     f"{issue['site']}: {issue['issue']}", font_size=9, color=WHITE)

        # アクション（右）
        add_text_box(slide, Inches(5.4), y, Inches(4.0), Inches(0.4),
                     f"→ {issue['action']}", font_size=9, color=WHITE)

    if len(issues) > max_display:
        add_text_box(slide, Inches(0.5), Inches(2.4 + max_display * 0.45), Inches(9), Inches(0.4),
                     f"他 {len(issues) - max_display} 件の問題あり", font_size=10, color=GRAY)


def main():
    base_dir = Path(__file__).resolve().parent.parent
    json_path = base_dir / "results" / "latest.json"

    if not json_path.exists():
        print("エラー: results/latest.json が見つかりません")
        print("先に python3 scripts/audit.py を実行してください")
        return

    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    timestamp = data["timestamp"]
    sites = data["sites"]
    date_str = timestamp.split(" ")[0] if " " in timestamp else timestamp

    print(f"パワポレポート生成中...")
    print(f"  データ: {json_path}")
    print(f"  日時: {timestamp}")
    print(f"  サイト数: {len(sites)}")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_cover_slide(prs, timestamp)
    create_summary_slide(prs, sites)
    create_detail_slide(prs, sites)
    create_seo_slide(prs, sites)
    create_issues_slide(prs, sites)

    reports_dir = base_dir / "reports"
    reports_dir.mkdir(parents=True, exist_ok=True)
    output_path = reports_dir / f"監査レポート_{date_str}.pptx"
    prs.save(str(output_path))
    print(f"\n  レポート出力: {output_path}")


if __name__ == "__main__":
    main()
